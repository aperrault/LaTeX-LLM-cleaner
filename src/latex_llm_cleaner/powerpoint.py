"""Extract text from PPTX files for LLM consumption."""

import sys
from pathlib import Path
from xml.etree.ElementTree import tostring as xml_tostring

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_OMML_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def extract_text_from_pptx(
    path: Path,
    verbose: bool = False,
    notes: bool = False,
    figure_summary_suffix: str = "_summary.txt",
    encoding: str = "utf-8",
) -> str:
    """Extract text from a PPTX as markdown."""
    prs = Presentation(str(path))
    base_dir = path.parent.resolve()
    pptx_stem = path.stem
    slides_md = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        md = _slide_to_markdown(
            slide, slide_num, notes, figure_summary_suffix, base_dir, encoding,
            verbose, pptx_stem,
        )
        slides_md.append(md)

    return "\n\n---\n\n".join(slides_md) + "\n"


def _slide_to_markdown(
    slide, slide_num, notes, figure_summary_suffix, base_dir, encoding, verbose,
    pptx_stem,
):
    """Convert a single slide to markdown."""
    # Title
    title_shape = slide.shapes.title
    if title_shape and title_shape.has_text_frame:
        title_text = title_shape.text_frame.text.strip()
    else:
        title_text = None

    if title_text:
        heading = f"# Slide {slide_num}: {title_text}"
    else:
        heading = f"# Slide {slide_num}"

    parts = [heading]

    # Collect IDs of shapes python-pptx already exposes, so we can
    # detect AlternateContent elements whose Choice branch is hidden.
    exposed_ids = set()
    for shape in slide.shapes:
        cNvPr = shape._element.find(f".//{{{_PML_NS}}}cNvPr")
        if cNvPr is not None:
            exposed_ids.add(cNvPr.get("id"))

    # Collect shapes (skip the title shape since we already used it)
    image_counter = 0
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        text, image_counter = _shape_to_text(
            shape, slide_num, image_counter, base_dir,
            figure_summary_suffix, encoding, verbose, pptx_stem,
        )
        if text:
            parts.append(text)

    # Process mc:AlternateContent elements whose Choice branch shapes
    # are not exposed by python-pptx (e.g. content with OMML math).
    sp_tree = slide._element.find(f"{{{_PML_NS}}}cSld/{{{_PML_NS}}}spTree")
    if sp_tree is not None:
        for ac in sp_tree.findall(f"{{{_MC_NS}}}AlternateContent"):
            choice = ac.find(f"{{{_MC_NS}}}Choice")
            if choice is None:
                continue
            for sp in choice:
                cNvPr = sp.find(f".//{{{_PML_NS}}}cNvPr")
                if cNvPr is not None and cNvPr.get("id") in exposed_ids:
                    continue  # already processed via python-pptx
                text = _extract_text_from_sp_element(sp)
                if text:
                    parts.append(text)

    # Speaker notes
    if notes and slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_text = notes_frame.text.strip()
        if notes_text:
            parts.append(f"> **Notes:** {notes_text}")

    return "\n\n".join(parts)


def _shape_to_text(shape, slide_num, image_counter, base_dir,
                   figure_summary_suffix, encoding, verbose, pptx_stem):
    """Convert a shape to text. Returns (text, updated_image_counter)."""
    # Group shape — recurse
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_parts = []
        for child in shape.shapes:
            text, image_counter = _shape_to_text(
                child, slide_num, image_counter, base_dir,
                figure_summary_suffix, encoding, verbose, pptx_stem,
            )
            if text:
                group_parts.append(text)
        return "\n\n".join(group_parts), image_counter

    # Table
    if shape.has_table:
        return _table_to_markdown(shape.table), image_counter

    # Picture / image (including placeholders with embedded images)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
        shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
    ):
        image_counter += 1
        summary = _find_image_summary(
            base_dir, pptx_stem, slide_num, image_counter,
            figure_summary_suffix, encoding,
        )
        if summary:
            return f"[Image: {summary}]", image_counter
        if verbose:
            print(
                f"Warning: no summary found for {pptx_stem}_slide{slide_num}_image{image_counter}",
                file=sys.stderr,
            )
        return "[Image]", image_counter

    # Text frame
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            text = _process_paragraph(para)
            if text:
                paragraphs.append(text)
        return "\n\n".join(paragraphs), image_counter

    return None, image_counter


def _extract_text_from_sp_element(sp):
    """Extract text from a raw sp XML element (for mc:AlternateContent shapes)."""
    txBody = sp.find(f"{{{_PML_NS}}}txBody")
    if txBody is None:
        return None
    paragraphs = []
    for p_elem in txBody.findall(f"{{{_DML_NS}}}p"):
        text = _process_paragraph_element(p_elem)
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs) if paragraphs else None


def _process_paragraph(paragraph):
    """Extract text from a python-pptx Paragraph object."""
    return _process_paragraph_element(paragraph._element)


def _process_paragraph_element(element):
    """Extract text from a paragraph XML element, preserving OMML math as XML."""
    ns_math = f"{{{_OMML_NS}}}"

    parts = []
    # Walk the paragraph's XML children to interleave text and math
    for child in element:
        tag = child.tag
        # OMML math paragraph (display math)
        if tag == f"{ns_math}oMathPara" or tag == f"{ns_math}oMath":
            math_xml = xml_tostring(child, encoding="unicode")
            parts.append(math_xml)
        # Regular text run (a:r in drawingML)
        elif tag.endswith("}r"):
            # Get text from <a:t> child
            for t_elem in child:
                if t_elem.tag.endswith("}t") and t_elem.text:
                    parts.append(t_elem.text)
        # a14:m wrapper around OMML math (used in mc:Choice branches)
        elif tag.endswith("}m"):
            for math_child in child:
                if ns_math in math_child.tag:
                    math_xml = xml_tostring(math_child, encoding="unicode")
                    parts.append(math_xml)

    result = "".join(parts).strip()
    return result if result else None


def _table_to_markdown(table):
    """Convert a pptx Table to markdown pipe-table format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

    if not rows:
        return ""

    # Insert separator after header row
    col_count = len(table.columns)
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    rows.insert(1, separator)

    return "\n".join(rows)


def _find_image_summary(base_dir, pptx_stem, slide_num, image_index, suffix, encoding):
    """Look for {pptx_stem}_slide{N}_image{M}{suffix} in base_dir."""
    stem = f"{pptx_stem}_slide{slide_num}_image{image_index}"
    summary_path = base_dir / (stem + suffix)
    if summary_path.is_file():
        return summary_path.read_text(encoding=encoding).strip()
    return None
