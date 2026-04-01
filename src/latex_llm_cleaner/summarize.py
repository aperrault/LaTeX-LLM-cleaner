"""Auto-generate figure summaries using the Gemini vision API."""

import mimetypes
import re
import sys
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from functools import wraps
from pathlib import Path

from google import genai
from google.genai import types

from .figures import _INCLUDEGRAPHICS_RE, _IMAGE_EXTENSIONS, _find_summary

_PROMPT = (
    "Describe this figure so that someone who cannot see it has all the same "
    "information. For charts and plots, extract the data into tables. For "
    "diagrams, describe the structure and relationships. For photographs, "
    "describe what is shown. Do not editorialize or interpret beyond what "
    "the figure itself conveys."
)

_MODEL = "gemini-3.1-flash-lite-preview"

_MAX_WORKERS = 4


def _retry_with_backoff(max_retries=3, base_delay=10, max_delay=60):
    """Retry decorator with exponential backoff for transient API errors."""

    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            last_exception = None
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    last_exception = e
                    error_str = str(e).lower()
                    retryable = any(
                        x in error_str
                        for x in [
                            "timeout",
                            "503",
                            "502",
                            "504",
                            "429",
                            "unavailable",
                            "cancelled",
                            "overloaded",
                            "rate limit",
                            "connection",
                        ]
                    )
                    if not retryable or attempt == max_retries - 1:
                        raise
                    delay = min(base_delay * (2**attempt), max_delay)
                    print(
                        f"  Retrying in {delay}s (attempt {attempt + 1}/{max_retries})...",
                        file=sys.stderr,
                    )
                    time.sleep(delay)
            raise last_exception  # type: ignore[misc]

        return wrapper

    return decorator


def _print_progress(done: int, total: int) -> None:
    """Print an in-place progress counter to stderr."""
    print(f"\r  Summarizing: {done}/{total} done", end="", file=sys.stderr, flush=True)


def _resolve_image_path(base_dir: Path, img_path_str: str) -> Path | None:
    """Find the actual image file on disk, probing extensions if needed."""
    img_path = base_dir / img_path_str

    if img_path.is_file():
        return img_path

    # If no extension, try common image extensions
    if img_path.suffix.lower() not in _IMAGE_EXTENSIONS:
        for ext in _IMAGE_EXTENSIONS:
            candidate = base_dir / (img_path_str + ext)
            if candidate.is_file():
                return candidate

    return None


def _get_mime_type(image_path: Path) -> str:
    """Detect MIME type from file extension."""
    mime, _ = mimetypes.guess_type(str(image_path))
    return mime or "application/octet-stream"


_GEMINI_SUPPORTED_MIMES = {"image/png", "image/jpeg", "image/webp", "image/heic", "image/heif"}


def _ensure_supported_format(image_bytes: bytes, mime_type: str) -> tuple[bytes, str]:
    """Convert unsupported image formats (GIF, TIFF, BMP, etc.) to PNG.

    Raises ValueError for formats that cannot be converted (e.g. WMF/EMF).
    """
    if mime_type in _GEMINI_SUPPORTED_MIMES:
        return image_bytes, mime_type
    import pymupdf

    try:
        pix = pymupdf.Pixmap(image_bytes)
        return pix.tobytes("png"), "image/png"
    except Exception:
        raise ValueError(f"unsupported image format: {mime_type}")


@_retry_with_backoff()
def _call_gemini_bytes(client, image_bytes: bytes, mime_type: str, prompt: str) -> str:
    """Send image bytes + prompt to Gemini, return generated text."""
    image_bytes, mime_type = _ensure_supported_format(image_bytes, mime_type)
    response = client.models.generate_content(
        model=_MODEL,
        contents=[
            types.Content(
                parts=[
                    types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                    types.Part.from_text(text=prompt),
                ]
            )
        ],
    )
    return response.text


def _call_gemini(client, image_path: Path, prompt: str) -> str:
    """Send image file + prompt to Gemini, return generated text."""
    return _call_gemini_bytes(
        client, image_path.read_bytes(), _get_mime_type(image_path), prompt
    )


def auto_summarize_figures(content: str, base_dir: Path, options: dict) -> str:
    """Generate summary files for figures that lack them.

    This function has side effects (writes _summary.txt files) but returns
    the content string unchanged.  The subsequent substitute_figures step
    will pick up the newly created summaries.
    """
    api_key = options.get("google_api_key")
    if not api_key:
        print(
            "Error: No Google API key found. Set GOOGLE_API_KEY or use --google-api-key.",
            file=sys.stderr,
        )
        sys.exit(1)

    verbose = options.get("verbose", False)
    suffix = options.get("figure_summary_suffix", "_summary.txt")
    encoding = options.get("encoding", "utf-8")

    # Collect unique image paths that need summaries
    seen: set[str] = set()
    work_items: list[tuple[str, Path, Path]] = []  # (label, image_path, summary_path)
    skipped = 0
    for m in _INCLUDEGRAPHICS_RE.finditer(content):
        img_path_str = m.group(1).strip()
        if img_path_str in seen:
            continue
        seen.add(img_path_str)

        existing = _find_summary(base_dir, img_path_str, suffix, encoding)
        if existing is not None:
            if verbose:
                print(f"  Skipping {img_path_str} (summary exists)", file=sys.stderr)
            skipped += 1
            continue

        image_path = _resolve_image_path(base_dir, img_path_str)
        if image_path is None:
            if verbose:
                print(
                    f"  Warning: image file not found for {img_path_str}",
                    file=sys.stderr,
                )
            continue

        if image_path.suffix.lower() in _IMAGE_EXTENSIONS:
            stem_path = image_path.with_suffix("")
        else:
            stem_path = image_path
        summary_path = Path(str(stem_path) + suffix)
        work_items.append((img_path_str, image_path, summary_path))

    if not work_items:
        if verbose:
            print("  No figures found to summarize.", file=sys.stderr)
        return content

    client = genai.Client(api_key=api_key)
    total = len(work_items)
    generated = 0

    _print_progress(0, total)

    def _do_one(item):
        label, image_path, summary_path = item
        summary_text = _call_gemini(client, image_path, _PROMPT)
        summary_path.write_text(summary_text, encoding=encoding)
        return label, summary_path

    with ThreadPoolExecutor(max_workers=_MAX_WORKERS) as pool:
        futures = {pool.submit(_do_one, item): item for item in work_items}
        for future in as_completed(futures):
            label, _, _ = futures[future]
            try:
                _, summary_path = future.result()
                generated += 1
                if verbose:
                    print(f"\r  Wrote {summary_path}" + " " * 20, file=sys.stderr)
            except Exception as e:
                print(
                    f"\r  Warning: API error for {label}: {e}" + " " * 20,
                    file=sys.stderr,
                )
            _print_progress(generated, total)

    print(file=sys.stderr)  # newline after progress
    if verbose:
        print(
            f"  Generated {generated} summary file(s) "
            f"({skipped} skipped).",
            file=sys.stderr,
        )

    return content


def auto_summarize_pptx(path: Path, options: dict) -> None:
    """Generate summary files for PPTX images that lack them.

    Writes {pptx_stem}_slide{N}_image{M}_summary.txt files next to the PPTX
    so the existing extraction pipeline picks them up.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    api_key = options.get("google_api_key")
    if not api_key:
        print(
            "Error: No Google API key found. Set GOOGLE_API_KEY or use --google-api-key.",
            file=sys.stderr,
        )
        sys.exit(1)

    verbose = options.get("verbose", False)
    suffix = options.get("figure_summary_suffix", "_summary.txt")
    encoding = options.get("encoding", "utf-8")
    base_dir = path.parent.resolve()
    pptx_stem = path.stem

    prs = Presentation(str(path))

    # First pass: collect all work items
    work_items: list[tuple[str, bytes, str, Path]] = []  # (stem, blob, mime, summary_path)
    skipped = 0

    def _collect_shape(shape, slide_num, image_counter):
        nonlocal skipped

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                image_counter = _collect_shape(child, slide_num, image_counter)
            return image_counter

        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return image_counter

        image_counter += 1
        stem = f"{pptx_stem}_slide{slide_num}_image{image_counter}"
        summary_path = base_dir / (stem + suffix)

        if summary_path.is_file():
            if verbose:
                print(f"  Skipping {stem} (summary exists)", file=sys.stderr)
            skipped += 1
            return image_counter

        work_items.append((
            stem,
            shape.image.blob,
            shape.image.content_type,
            summary_path,
        ))
        return image_counter

    for slide_num, slide in enumerate(prs.slides, start=1):
        image_counter = 0
        for shape in slide.shapes:
            image_counter = _collect_shape(shape, slide_num, image_counter)

    if not work_items:
        if verbose:
            print("  No images found to summarize.", file=sys.stderr)
        return

    _run_batch_summarize(work_items, api_key, encoding, verbose, skipped)


# Minimum dimensions for a picture marker to be considered a real figure
_MIN_FIGURE_DIM = 64

_PICTURE_MARKER_RE = re.compile(
    r"\*\*==> picture \[(\d+) x (\d+)\] intentionally omitted <==\*\*"
)


def auto_summarize_pdf(path: Path, options: dict) -> None:
    """Generate summary files for significant figures in a PDF.

    Extracts embedded images and cropped picture regions, sends them to
    Gemini for summarization, and writes {pdf_stem}_page{N}_image{M}_summary.txt
    files next to the PDF.
    """
    import fitz
    import pymupdf4llm

    api_key = options.get("google_api_key")
    if not api_key:
        print(
            "Error: No Google API key found. Set GOOGLE_API_KEY or use --google-api-key.",
            file=sys.stderr,
        )
        sys.exit(1)

    verbose = options.get("verbose", False)
    suffix = options.get("figure_summary_suffix", "_summary.txt")
    encoding = options.get("encoding", "utf-8")
    base_dir = path.parent.resolve()
    pdf_stem = path.stem

    doc = fitz.open(path)
    chunks = pymupdf4llm.to_markdown(str(path), page_chunks=True)

    # Collect work items: (stem, image_bytes, mime_type, summary_path)
    work_items: list[tuple[str, bytes, str, Path]] = []
    skipped = 0

    for page_num, chunk in enumerate(chunks, start=1):
        page = doc[page_num - 1]

        # Count significant picture markers on this page
        markers = _PICTURE_MARKER_RE.findall(chunk["text"])
        significant_markers = [
            (int(w), int(h)) for w, h in markers
            if int(w) > _MIN_FIGURE_DIM and int(h) > _MIN_FIGURE_DIM
        ]

        if not significant_markers:
            continue

        # Strategy: try embedded images first, fall back to picture box crops
        embedded = []
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            base = doc.extract_image(xref)
            if base["width"] > _MIN_FIGURE_DIM and base["height"] > _MIN_FIGURE_DIM:
                ext = base["ext"]
                mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                embedded.append((base["image"], mime))

        # Get picture boxes for cropped rendering fallback
        pic_boxes = [b for b in chunk["page_boxes"] if b["class"] == "picture"]

        # Build image list: prefer embedded, supplement with cropped boxes
        page_images: list[tuple[bytes, str]] = []
        if embedded:
            page_images.extend(embedded)
        # If we have more significant markers than embedded images, use box crops
        remaining = len(significant_markers) - len(embedded)
        if remaining > 0 and pic_boxes:
            for box in pic_boxes[:remaining]:
                bbox = fitz.Rect(box["bbox"])
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(clip=bbox, matrix=mat)
                page_images.append((pix.tobytes("png"), "image/png"))

        # Create work items for each image on this page
        for img_idx, (img_bytes, mime_type) in enumerate(page_images, start=1):
            stem = f"{pdf_stem}_page{page_num}_image{img_idx}"
            summary_path = base_dir / (stem + suffix)

            if summary_path.is_file():
                if verbose:
                    print(f"  Skipping {stem} (summary exists)", file=sys.stderr)
                skipped += 1
                continue

            work_items.append((stem, img_bytes, mime_type, summary_path))

    doc.close()

    if not work_items:
        if verbose:
            print("  No figures found to summarize.", file=sys.stderr)
        return

    _run_batch_summarize(work_items, api_key, encoding, verbose, skipped)


def _run_batch_summarize(
    work_items: list[tuple[str, bytes, str, Path]],
    api_key: str,
    encoding: str,
    verbose: bool,
    skipped: int,
) -> None:
    """Run Gemini summarization on a batch of (stem, bytes, mime, path) items."""
    client = genai.Client(api_key=api_key)
    total = len(work_items)
    generated = 0

    _print_progress(0, total)

    def _do_one(item):
        stem, image_bytes, content_type, summary_path = item
        summary_text = _call_gemini_bytes(client, image_bytes, content_type, _PROMPT)
        summary_path.write_text(summary_text, encoding=encoding)
        return stem, summary_path

    with ThreadPoolExecutor(max_workers=_MAX_WORKERS) as pool:
        futures = {pool.submit(_do_one, item): item for item in work_items}
        for future in as_completed(futures):
            stem = futures[future][0]
            try:
                _, summary_path = future.result()
                generated += 1
                if verbose:
                    print(f"\r  Wrote {summary_path}" + " " * 20, file=sys.stderr)
            except Exception as e:
                print(
                    f"\r  Warning: API error for {stem}: {e}" + " " * 20,
                    file=sys.stderr,
                )
            _print_progress(generated, total)

    print(file=sys.stderr)  # newline after progress
    if verbose:
        total_images = generated + skipped
        print(
            f"  Generated {generated} summary file(s) "
            f"({skipped} skipped, {total_images} total images).",
            file=sys.stderr,
        )


def auto_summarize_docx(path: Path, options: dict) -> None:
    """Generate summary files for DOCX images that lack them.

    Writes {docx_stem}_image{M}_summary.txt files next to the DOCX
    so the existing extraction pipeline picks them up.
    """
    from lxml import etree

    from docx import Document

    _DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    _WP_NS = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
    _REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    _WML_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    api_key = options.get("google_api_key")
    if not api_key:
        print(
            "Error: No Google API key found. Set GOOGLE_API_KEY or use --google-api-key.",
            file=sys.stderr,
        )
        sys.exit(1)

    verbose = options.get("verbose", False)
    suffix = options.get("figure_summary_suffix", "_summary.txt")
    encoding = options.get("encoding", "utf-8")
    base_dir = path.parent.resolve()
    docx_stem = path.stem

    doc = Document(str(path))

    # Walk body elements in the same order as docx.py extraction.
    # Only count inline drawings (wp:inline), not anchor drawings
    # (text boxes, floating decorations).
    work_items: list[tuple[str, bytes, str, Path]] = []
    skipped = 0
    image_counter = 0

    for child in doc.element.body:
        tag = etree.QName(child.tag).localname
        if tag != "p":
            continue

        for run in child.findall(f".//{{{_WML_NS}}}r"):
            # Use descendant search since drawings may be inside
            # mc:AlternateContent/mc:Choice wrappers.
            for inline in run.findall(f".//{{{_WP_NS}}}inline"):
                for blip in inline.findall(f".//{{{_DML_NS}}}blip"):
                    image_counter += 1
                    stem = f"{docx_stem}_image{image_counter}"
                    summary_path = base_dir / (stem + suffix)

                    if summary_path.is_file():
                        if verbose:
                            print(
                                f"  Skipping {stem} (summary exists)",
                                file=sys.stderr,
                            )
                        skipped += 1
                        continue

                    rId = blip.get(f"{{{_REL_NS}}}embed")
                    if not rId:
                        continue

                    try:
                        image_part = doc.part.rels[rId].target_part
                        work_items.append((
                            stem,
                            image_part.blob,
                            image_part.content_type,
                            summary_path,
                        ))
                    except (KeyError, AttributeError) as e:
                        if verbose:
                            print(
                                f"  Warning: could not extract image "
                                f"for {stem}: {e}",
                                file=sys.stderr,
                            )

    if not work_items:
        if verbose:
            print("  No images found to summarize.", file=sys.stderr)
        return

    _run_batch_summarize(work_items, api_key, encoding, verbose, skipped)
