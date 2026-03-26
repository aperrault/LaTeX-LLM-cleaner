"""Tests for PPTX text extraction."""

import pytest
from pathlib import Path
from unittest.mock import patch

from latex_llm_cleaner.powerpoint import extract_text_from_pptx


try:
    import pptx

    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

skip_no_pptx = pytest.mark.skipif(
    not HAS_PYTHON_PPTX, reason="python-pptx not installed"
)


def _create_test_pptx(path, slides_data):
    """Create a simple PPTX for testing.

    slides_data: list of dicts with optional keys:
        - title: str
        - body: str or list[str]
        - notes: str
        - table: list[list[str]]  (rows of cells)
        - image_path: Path to an image file to embed
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for slide_info in slides_data:
        if "table" in slide_info:
            layout = prs.slide_layouts[5]  # blank
            slide = prs.slides.add_slide(layout)
            rows_data = slide_info["table"]
            n_rows = len(rows_data)
            n_cols = len(rows_data[0]) if rows_data else 0
            table_shape = slide.shapes.add_table(
                n_rows, n_cols, Inches(1), Inches(2), Inches(6), Inches(3)
            )
            for r, row_data in enumerate(rows_data):
                for c, cell_text in enumerate(row_data):
                    table_shape.table.cell(r, c).text = cell_text
        elif "title" in slide_info or "body" in slide_info:
            layout = prs.slide_layouts[1]  # title + content
            slide = prs.slides.add_slide(layout)
            if "title" in slide_info and slide.shapes.title:
                slide.shapes.title.text = slide_info["title"]
            if "body" in slide_info:
                body_ph = slide.placeholders[1]
                body = slide_info["body"]
                if isinstance(body, str):
                    body_ph.text = body
                else:
                    body_ph.text = body[0]
                    for line in body[1:]:
                        body_ph.text_frame.add_paragraph().text = line
        else:
            layout = prs.slide_layouts[5]  # blank
            slide = prs.slides.add_slide(layout)

        if "image_path" in slide_info:
            from pptx.util import Inches

            slide.shapes.add_picture(
                str(slide_info["image_path"]), Inches(1), Inches(1), Inches(2), Inches(2)
            )

        if "notes" in slide_info:
            slide.notes_slide.notes_text_frame.text = slide_info["notes"]

    prs.save(str(path))


@skip_no_pptx
def test_basic_extraction(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [{"title": "Hello", "body": "World"}])
    result = extract_text_from_pptx(pptx_path)
    assert "Hello" in result
    assert "World" in result


@skip_no_pptx
def test_slide_numbering(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"title": "First", "body": "A"},
        {"title": "Second", "body": "B"},
    ])
    result = extract_text_from_pptx(pptx_path)
    assert "# Slide 1: First" in result
    assert "# Slide 2: Second" in result


@skip_no_pptx
def test_slide_separator(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"title": "One", "body": "X"},
        {"title": "Two", "body": "Y"},
    ])
    result = extract_text_from_pptx(pptx_path)
    assert "\n\n---\n\n" in result


@skip_no_pptx
def test_no_title(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [{}])  # blank slide
    result = extract_text_from_pptx(pptx_path)
    assert "# Slide 1" in result


@skip_no_pptx
def test_notes_off_by_default(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"title": "Slide", "body": "Content", "notes": "Secret notes"}
    ])
    result = extract_text_from_pptx(pptx_path)
    assert "Secret notes" not in result


@skip_no_pptx
def test_notes_included_when_enabled(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"title": "Slide", "body": "Content", "notes": "Speaker notes here"}
    ])
    result = extract_text_from_pptx(pptx_path, notes=True)
    assert "Speaker notes here" in result
    assert "> **Notes:**" in result


@skip_no_pptx
def test_table_extraction(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"table": [["Name", "Score"], ["Alice", "95"], ["Bob", "87"]]}
    ])
    result = extract_text_from_pptx(pptx_path)
    assert "| Name | Score |" in result
    assert "| --- | --- |" in result
    assert "| Alice | 95 |" in result
    assert "| Bob | 87 |" in result


@skip_no_pptx
def test_image_summary_lookup(tmp_path):
    """Image with a matching summary file should be replaced."""
    # Create a tiny PNG for embedding
    img_path = tmp_path / "tiny.png"
    _create_tiny_png(img_path)

    # Create summary file
    summary_path = tmp_path / "slide1_image1_summary.txt"
    summary_path.write_text("A bar chart showing results.")

    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"title": "Results", "image_path": img_path}
    ])
    result = extract_text_from_pptx(pptx_path)
    assert "[Image: A bar chart showing results.]" in result


@skip_no_pptx
def test_image_no_summary(tmp_path):
    """Image without summary should produce [Image] placeholder."""
    img_path = tmp_path / "tiny.png"
    _create_tiny_png(img_path)

    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path, [
        {"title": "Fig", "image_path": img_path}
    ])
    result = extract_text_from_pptx(pptx_path)
    assert "[Image]" in result


def test_import_error_message(tmp_path):
    """When python-pptx is not installed, should exit with helpful message."""
    pptx_path = tmp_path / "test.pptx"
    pptx_path.write_bytes(b"fake")

    with patch.dict("sys.modules", {"pptx": None}):
        with pytest.raises(SystemExit):
            extract_text_from_pptx(pptx_path)


def _create_tiny_png(path):
    """Create a minimal valid PNG file."""
    import struct
    import zlib

    # 1x1 white PNG
    def chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw = zlib.compress(b"\x00\xff\xff\xff")
    idat = chunk(b"IDAT", raw)
    iend = chunk(b"IEND", b"")
    path.write_bytes(sig + ihdr + idat + iend)
